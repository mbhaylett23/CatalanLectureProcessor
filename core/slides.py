"""PowerPoint slide generation from structured summary data."""

import os
import logging
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from core.config import SLIDE_MAX_BULLETS

logger = logging.getLogger(__name__)

# Colour palette
COLOR_TITLE = RGBColor(0x1A, 0x47, 0x8A)  # Dark blue
COLOR_BODY = RGBColor(0x33, 0x33, 0x33)   # Dark grey
COLOR_ACCENT = RGBColor(0x2E, 0x86, 0xC1)  # Medium blue
COLOR_LIGHT = RGBColor(0x85, 0x92, 0x9E)   # Light grey


class SlideGenerator:
    """Generate PowerPoint presentations from lecture summary data."""

    def generate(self, summary_data: dict, title: str = "Lecture Summary",
                 output_path: str = None) -> str:
        """Create a PowerPoint presentation from summary data.

        Args:
            summary_data: Dict with main_topics, sections, key_terms, detailed_summary.
            title: Lecture title for the title slide.
            output_path: Where to save. Auto-generated if None.

        Returns:
            Path to the generated .pptx file.
        """
        prs = Presentation()

        # Set 16:9 widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        self._add_title_slide(prs, title, f"Generated {date.today().isoformat()}")

        # Overview slide with main topics
        main_topics = summary_data.get("main_topics", [])
        if main_topics:
            self._add_content_slide(prs, "Overview", main_topics)

        # Content slides from sections
        sections = summary_data.get("sections", [])
        for section in sections:
            section_title = section.get("title", "")
            bullets = section.get("bullets", [])
            if section_title and bullets:
                self._add_content_slide(prs, section_title, bullets)

        # Detailed summary slide (if available)
        detailed = summary_data.get("detailed_summary", "")
        if detailed:
            self._add_text_slide(prs, "Detailed Summary", detailed)

        # Key terms slide
        key_terms = summary_data.get("key_terms", [])
        if key_terms:
            self._add_content_slide(prs, "Key Terms & Concepts", key_terms)

        # End slide
        self._add_end_slide(prs)

        # Save
        if output_path is None:
            output_path = os.path.join(
                os.getcwd(), f"lecture_slides_{date.today().isoformat()}.pptx"
            )
        prs.save(output_path)
        logger.info("Slides saved to: %s", output_path)
        return output_path

    def _add_title_slide(self, prs: Presentation, title_text: str, subtitle_text: str):
        """Add a title slide."""
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        title = slide.shapes.title
        title.text = title_text
        for para in title.text_frame.paragraphs:
            para.font.size = Pt(36)
            para.font.color.rgb = COLOR_TITLE
            para.font.bold = True

        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = subtitle_text
            for para in subtitle.text_frame.paragraphs:
                para.font.size = Pt(18)
                para.font.color.rgb = COLOR_LIGHT

    def _add_content_slide(self, prs: Presentation, title_text: str, bullets: list[str]):
        """Add a content slide with bullet points. Overflows to continuation slides."""
        # Split bullets into groups of SLIDE_MAX_BULLETS
        for page_idx in range(0, len(bullets), SLIDE_MAX_BULLETS):
            page_bullets = bullets[page_idx:page_idx + SLIDE_MAX_BULLETS]
            is_continuation = page_idx > 0
            slide_title = f"{title_text} (cont.)" if is_continuation else title_text

            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)

            title = slide.shapes.title
            title.text = slide_title
            for para in title.text_frame.paragraphs:
                para.font.size = Pt(28)
                para.font.color.rgb = COLOR_TITLE
                para.font.bold = True

            if len(slide.placeholders) > 1:
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.clear()

                for i, bullet in enumerate(page_bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = bullet
                    p.font.size = Pt(18)
                    p.font.color.rgb = COLOR_BODY
                    p.space_after = Pt(8)
                    p.level = 0

    def _add_text_slide(self, prs: Presentation, title_text: str, body_text: str):
        """Add a slide with paragraph text (for detailed summary)."""
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        title = slide.shapes.title
        title.text = title_text
        for para in title.text_frame.paragraphs:
            para.font.size = Pt(28)
            para.font.color.rgb = COLOR_TITLE
            para.font.bold = True

        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            tf.word_wrap = True

            paragraphs = body_text.split("\n\n")
            for i, para_text in enumerate(paragraphs):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = para_text.strip()
                p.font.size = Pt(14)
                p.font.color.rgb = COLOR_BODY
                p.space_after = Pt(12)

    def _add_end_slide(self, prs: Presentation):
        """Add an end slide."""
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        title = slide.shapes.title
        title.text = "End of Summary"
        for para in title.text_frame.paragraphs:
            para.font.size = Pt(36)
            para.font.color.rgb = COLOR_TITLE
            para.alignment = PP_ALIGN.CENTER

        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = "Generated with Catalan Lecture Processor"
            for para in subtitle.text_frame.paragraphs:
                para.font.size = Pt(16)
                para.font.color.rgb = COLOR_LIGHT
                para.alignment = PP_ALIGN.CENTER
