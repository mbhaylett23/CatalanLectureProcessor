"""Tests for the PowerPoint slide generator."""

import sys
import os
import tempfile

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from core.slides import SlideGenerator


def test_creates_valid_pptx():
    gen = SlideGenerator()
    summary = {
        "main_topics": ["Topic A", "Topic B", "Topic C"],
        "sections": [
            {"title": "Section 1", "bullets": ["Point 1", "Point 2"]},
            {"title": "Section 2", "bullets": ["Point 3", "Point 4"]},
        ],
        "key_terms": ["Term 1", "Term 2"],
        "detailed_summary": "This is a test summary paragraph.",
    }
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name

    try:
        result_path = gen.generate(summary, title="Test Lecture", output_path=path)
        assert os.path.exists(result_path)
        assert result_path.endswith(".pptx")
        assert os.path.getsize(result_path) > 0

        # Verify it can be opened
        from pptx import Presentation
        prs = Presentation(result_path)
        # Title + Overview + 2 sections + Detailed + Key Terms + End = 7 slides
        assert len(prs.slides) >= 5
    finally:
        os.unlink(path)


def test_overflow_creates_continuation():
    gen = SlideGenerator()
    summary = {
        "main_topics": [],
        "sections": [
            {"title": "Big Section", "bullets": [f"Point {i}" for i in range(15)]},
        ],
        "key_terms": [],
    }
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name

    try:
        gen.generate(summary, output_path=path)
        from pptx import Presentation
        prs = Presentation(path)
        # 15 bullets / 6 per slide = 3 content slides + title + end = 5
        titles = [s.shapes.title.text for s in prs.slides]
        assert any("cont." in t for t in titles)
    finally:
        os.unlink(path)


def test_empty_summary():
    gen = SlideGenerator()
    summary = {
        "main_topics": [],
        "sections": [],
        "key_terms": [],
    }
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name

    try:
        gen.generate(summary, output_path=path)
        from pptx import Presentation
        prs = Presentation(path)
        # Title + End = 2 slides minimum
        assert len(prs.slides) >= 2
    finally:
        os.unlink(path)


if __name__ == "__main__":
    test_creates_valid_pptx()
    test_overflow_creates_continuation()
    test_empty_summary()
    print("All slides tests passed!")
